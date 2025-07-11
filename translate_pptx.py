from pptx import Presentation
from pptx.util import Inches
from googletrans import Translator, constants
import asyncio

async def translate_text(text, lang):
    translator = Translator()
    translation = await translator.translate(text, dest=lang)
    return translation.text

# Load the PowerPoint file
path='c:/vsprojects/coldfusion scripts/documents/'
ppt_file = 'Sample ClipTraining Lesson.pptx'
filename=ppt_file.partition('.')[0]
outfilename=filename+"_fr.pptx"
inppt = Presentation(path + ppt_file) # open the ppt file
outppt = Presentation() # open an unnamed powerpoint file

# Iterate through the slide layouts to get their names
for layout in inppt.slide_layouts:
    print(f"Layout Name: {layout.name}")
    print("\n" + "-"*40 + "\n") # print a separator line Iterate through each slide
    
for slide_number, slide in enumerate(inppt.slides, start=1):
    print(f"Slide {slide_number}:")
    # read the title and text of the slide and write them to the new slide
    slide_title = slide.shapes.title.text if slide.shapes.title else "No Title"
    slide_text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])

    #title = asyncio.run(translate_text(slide.shapes.title, 'fr'))
    #title = asyncio.run(translate_text(slide.shapes.title, 'fr'))
    print(f"Title: {slide_title}")
    print(f"Text: {slide_text}")
    print("\n" + "-"*40 + "\n") # print a separator line

# outppt.save('./documents/translated_presentation.pptx')